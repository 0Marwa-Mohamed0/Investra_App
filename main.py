from fastapi import FastAPI, UploadFile, File, Form
from fastapi.middleware.cors import CORSMiddleware
from pptx import Presentation
from openai import OpenAI
from supabase import create_client
from dotenv import load_dotenv
from docx import Document
import PyPDF2
import pandas as pd
import os
import io
import json
import uuid
from datetime import datetime, timezone

app = FastAPI()
load_dotenv()

OPENROUTER_API_KEY = os.getenv("OPENROUTER_API_KEY")
SUPABASE_URL = os.getenv("SUPABASE_URL")
SUPABASE_KEY = os.getenv("SUPABASE_KEY")

supabase = create_client(SUPABASE_URL, SUPABASE_KEY)

client = OpenAI(
    base_url="https://openrouter.ai/api/v1",
    api_key=OPENROUTER_API_KEY
)

MODEL_NAME = "openai/gpt-4o-mini"

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

# ============================================================
# دالة تنظيف UUID
# ============================================================
def clean_uuid(value):
    if not value:
        return None
    cleaned = str(value).strip()
    if cleaned in ["", "0", "string", "null", "none", "undefined"]:
        return None
    return cleaned

# ============================================================
# دالة توليد عنوان ذكي للمحادثة
# ============================================================
def generate_smart_title(user_message, ai_response):
    prompt = f"""
    Create a very short title (max 5 words) for this chat conversation based on the following interaction. 
    Return ONLY the title, nothing else.
    
    User: {user_message}
    AI: {ai_response}
    """
    try:
        response = client.chat.completions.create(
            model=MODEL_NAME,
            messages=[{"role": "user", "content": prompt}],
            max_tokens=20
        )
        return response.choices[0].message.content.strip().replace('"', '')
    except:
        return user_message[:30] # Fallback في حال فشل الـ AI

# ============================================================
# دوال استخراج النص من الملفات
# ============================================================
def extract_text_from_pptx(file_bytes):
    prs = Presentation(io.BytesIO(file_bytes))
    full_text = []
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text.strip())
        if slide_texts:
            full_text.append(f"[Slide {slide_num}]\n" + "\n".join(slide_texts))
    return "\n\n".join(full_text)

def extract_text_from_docx(file_bytes):
    doc = Document(io.BytesIO(file_bytes))
    return "\n".join([para.text for para in doc.paragraphs])

def extract_text_from_pdf(file_bytes):
    pdf_reader = PyPDF2.PdfReader(io.BytesIO(file_bytes))
    text = ""
    for page in pdf_reader.pages:
        page_text = page.extract_text()
        if page_text:
            text += page_text + "\n"
    return text

def extract_text_from_excel(file_bytes):
    df = pd.read_excel(io.BytesIO(file_bytes))
    return df.to_string()

def build_evaluation_prompt(extracted_text):
    return f"""
You are an expert business analyst and investment advisor evaluating a startup pitch deck.

Analyze the following content and provide:
1. Overall Rating (out of 5) -> Format your first line EXACTLY like this: "Rating: X.X" where X.X is the number.
2. Business Analytics Score (out of 5) with explanation
3. Marketing Plan Score (out of 5) with explanation
4. feasibility study Score (out of 5) with explanation
5. business Plan Score (out of 5) with explanation
6. Strengths (3-5 points)
7. Weaknesses (3-5 points)
8. Investor Attractiveness (High / Medium / Low) with reason
9. Recommendations for improvement

Be specific, professional, and data-driven.

--- CONTENT TO EVALUATE ---
{extracted_text}
---

Respond in the same language the content is written in.
"""

# ============================================================
# /evaluate — تقييم ملف وحفظ النتيجة في Supabase
# ============================================================
@app.post("/evaluate")
async def evaluate_pitch(
    file: UploadFile = File(...),
    idea_id: str = Form(None),
    user_id: str = Form(None)
):
    idea_id = clean_uuid(idea_id)
    user_id = clean_uuid(user_id)

    file_bytes = await file.read()
    filename = file.filename.lower()
    extracted_text = ""

    if filename.endswith(".pptx"):
        extracted_text = extract_text_from_pptx(file_bytes)
    elif filename.endswith(".docx"):
        extracted_text = extract_text_from_docx(file_bytes)
    elif filename.endswith(".pdf"):
        extracted_text = extract_text_from_pdf(file_bytes)
    elif filename.endswith((".xlsx", ".xls")):
        extracted_text = extract_text_from_excel(file_bytes)
    else:
        return {"error": "نوع الملف غير مدعوم. الأنواع المدعومة: .pptx, .docx, .pdf, .xlsx"}

    if not extracted_text.strip():
        return {"error": "Could not extract text from the file."}

    response = client.chat.completions.create(
        model=MODEL_NAME,
        messages=[
            {"role": "system", "content": "You are a professional startup evaluator."},
            {"role": "user", "content": build_evaluation_prompt(extracted_text)}
        ],
        max_tokens=2000,
    )

    result = response.choices[0].message.content
    current_time = datetime.now(timezone.utc).isoformat()

    session_data = supabase.table("AI_Sessions").insert({
        "session_id": str(uuid.uuid4()),
        "user_id": user_id,
        "title": f"تقييم: {filename}",
        "last_message_snippet": result[:100],
        "created_at": current_time,
        "updated_at": current_time
    }).execute()

    session_id = session_data.data[0]["session_id"]

    supabase.table("AI_Messages").insert({
        "message_id": str(uuid.uuid4()),
        "session_id": session_id,
        "sender_role": "user",
        "content": f"طلب تقييم الملف: {filename}",
        "file_name": filename,
        "created_at": current_time
    }).execute()

    supabase.table("AI_Messages").insert({
        "message_id": str(uuid.uuid4()),
        "session_id": session_id,
        "sender_role": "assistant",
        "content": result,
        "created_at": current_time
    }).execute()

    if idea_id:
        try:
            lines = result.strip().split('\n')
            rating_line = lines[0] if "ating" in lines[0] else lines[1]
            rating_num = float(''.join(c for c in rating_line if c.isdigit() or c == '.'))
            if 0 <= rating_num <= 5:
                supabase.table("ideas").update({
                    "ai_rating": rating_num
                }).eq("id", idea_id).execute()
        except:
            pass

    slides_count = extracted_text.count("[Slide") if filename.endswith(".pptx") else 0

    return {
        "filename": filename,
        "evaluation": result,
        "session_id": session_id,
        "content_length": len(extracted_text),
        "slides_analyzed": slides_count if slides_count > 0 else None,
        "status": "success"
    }

# ============================================================
# /chat — محادثة عادية مع حفظ في Supabase
# ============================================================
@app.post("/chat")
async def chat(
    message: str = Form(...),
    history: str = Form(default="[]"),
    session_id: str = Form(default=None),
    user_id: str = Form(default=None)
):
    session_id = clean_uuid(session_id)
    user_id = clean_uuid(user_id)

    try:
        chat_history = json.loads(history) if history and history.strip() not in ["", "[]"] else []
    except:
        chat_history = []

    messages = [
        {
            "role": "system",
            "content": """You are an AI assistant specialized in business analysis, 
            startup evaluation, and connecting entrepreneurs with investors. 
            Help users improve their business ideas and pitch decks."""
        }
    ]

    for msg in chat_history[-10:]:
        messages.append(msg)

    messages.append({"role": "user", "content": message})

    response = client.chat.completions.create(
        model=MODEL_NAME,
        messages=messages,
        max_tokens=1000,
    )

    ai_response = response.choices[0].message.content
    current_time = datetime.now(timezone.utc).isoformat()

    # 1. تحديث الـ Session
    if not session_id:
        # هنا يتم توليد العنوان الذكي عند بداية المحادثة
        smart_title = generate_smart_title(message, ai_response)
        session_data = supabase.table("AI_Sessions").insert({
            "session_id": str(uuid.uuid4()),
            "user_id": user_id,
            "title": smart_title, 
            "last_message_snippet": ai_response[:100],
        }).execute()
        session_id = session_data.data[0]["session_id"]
    else:
        supabase.table("AI_Sessions").update({
            "last_message_snippet": ai_response[:100],
            "updated_at": current_time
        }).eq("session_id", session_id).execute()

    # 2. حفظ رد الـ AI فقط في AI_Messages
    supabase.table("AI_Messages").insert({
        "message_id": str(uuid.uuid4()),
        "session_id": session_id,
        "sender_role": "assistant",
        "content": ai_response,
        "created_at": current_time
    }).execute()

    return {
        "response": ai_response,
        "session_id": session_id
    }

# ============================================================
# /rate-ideas
# ============================================================
def get_ideas():
    response = supabase.table("ideas").select("*").execute()
    return response.data

def update_rating(idea_id, rating):
    try:
        rating_num = float(''.join(c for c in str(rating) if c.isdigit() or c == '.'))
    except:
        rating_num = 0.0
    supabase.table("ideas").update({
        "ai_rating": rating_num
    }).eq("id", idea_id).execute()

def evaluate_idea(description, idea_docs):
    prompt = f"""
Rate this startup idea from 5 based on its description and documents:
Description: {description}
Documents: {idea_docs}
Return only a number between 0 and 5.
"""
    response = client.chat.completions.create(
        model=MODEL_NAME,
        messages=[{"role": "user", "content": prompt}],
    )
    return response.choices[0].message.content

@app.get("/rate-ideas")
def rate_ideas():
    ideas = get_ideas()
    for idea in ideas:
        description = idea.get("description", "")
        docs = idea.get("idea_docs", "")
        rating = evaluate_idea(description, docs)
        update_rating(idea["id"], rating)
    return {"status": "done", "rated": len(ideas)}

# ============================================================
# /test-read
# ============================================================
@app.post("/test-read")
async def test_read(file: UploadFile = File(...)):
    file_bytes = await file.read()
    filename = file.filename.lower()

    if filename.endswith(".pptx"):
        extracted_text = extract_text_from_pptx(file_bytes)
    elif filename.endswith(".docx"):
        extracted_text = extract_text_from_docx(file_bytes)
    elif filename.endswith(".pdf"):
        extracted_text = extract_text_from_pdf(file_bytes)
    elif filename.endswith((".xlsx", ".xls")):
        extracted_text = extract_text_from_excel(file_bytes)
    else:
        return {"error": "نوع الملف غير مدعوم."}

    return {
        "filename": filename,
        "content_preview": extracted_text[:500],
        "content_length": len(extracted_text),
        "status": "success"
    }

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)
